"""Offline document ingestion with provenance, chunking and idempotent upserts.

The default path is dependency-light and deterministic. Install optional PDF/PPT
parsers and replace the embedding hook with the configured model for production.
"""

import argparse
import hashlib
import json
import re
from pathlib import Path
from typing import Iterable

from app.config import get_settings


def _read_pages(path: Path) -> Iterable[tuple[int, str]]:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".markdown"}:
        yield 1, path.read_text(encoding="utf-8", errors="ignore")
        return
    if suffix == ".pdf":
        try:
            import fitz  # type: ignore

            document = fitz.open(path)
            for index, page in enumerate(document, start=1):
                yield index, page.get_text("text")
            return
        except Exception:
            pass
    if suffix == ".pptx":
        try:
            from pptx import Presentation  # type: ignore

            presentation = Presentation(path)
            for index, slide in enumerate(presentation.slides, start=1):
                text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                yield index, text
            return
        except Exception:
            pass
    # Keep an explicit fallback rather than silently discarding an unsupported file.
    yield 1, f"文件 {path.name} 尚未安装解析器，请使用视觉/OCR provider 补充内容。"


def _chunks(text: str, max_chars: int = 512) -> Iterable[str]:
    max_chars = max(128, min(8192, int(max_chars)))
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", text) if part.strip()]
    current = ""
    for paragraph in paragraphs:
        if len(paragraph) > max_chars:
            if current:
                yield current
                current = ""
            # Preserve every character from a long paragraph instead of
            # silently dropping everything after the first chunk.
            for start in range(0, len(paragraph), max_chars):
                yield paragraph[start : start + max_chars]
        elif len(current) + len(paragraph) + 1 <= max_chars:
            current = f"{current}\n{paragraph}".strip()
        else:
            if current:
                yield current
            current = paragraph[:max_chars]
    if current:
        yield current


def _hyde_questions(text: str, count: int) -> list[str]:
    if count <= 0:
        return []
    topic = re.sub(r"\s+", " ", text)[:80]
    templates = [f"如何理解{topic}？", f"{topic}的使用条件是什么？", f"如何应用{topic}解决题目？"]
    return templates[:count]


def _hash_embedding(text: str, dimensions: int = 32) -> list[float]:
    digest = hashlib.sha256(text.encode("utf-8")).digest()
    return [((digest[index % len(digest)] / 255.0) * 2.0) - 1.0 for index in range(dimensions)]


def build_rows(directory: Path, hyde_count: int = 0, chunk_chars: int = 512) -> list[dict]:
    rows: list[dict] = []
    for path in sorted(directory.rglob("*")):
        if not path.is_file() or path.suffix.lower() not in {".pdf", ".pptx", ".txt", ".md", ".markdown"}:
            continue
        # Derive the source id from stable logical path + file content, not the
        # temporary upload directory. Re-uploading the same document therefore
        # remains idempotent across browser indexing runs.
        relative_name = path.relative_to(directory).as_posix()
        file_digest = hashlib.sha256(path.read_bytes()).hexdigest()
        source_id = hashlib.sha256(f"{relative_name}\0{file_digest}".encode("utf-8")).hexdigest()[:16]
        for page, text in _read_pages(path):
            for chunk_index, chunk in enumerate(_chunks(text, chunk_chars), start=1):
                content_hash = hashlib.sha256(chunk.encode("utf-8")).hexdigest()
                chunk_id = f"{source_id}_p{page}_c{chunk_index}_{content_hash[:8]}"
                metadata = {
                    "source_id": source_id,
                    "chunk_id": chunk_id,
                    "filename": path.name,
                    "page": page,
                    "chapter": path.parent.name,
                    "chunk_index": chunk_index,
                    "content_hash": content_hash,
                    "parser_version": "1.1",
                }
                rows.append({"text": chunk, "metadata": metadata, "hyde_questions": _hyde_questions(chunk, hyde_count)})
    return rows


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--dir", required=True, type=Path)
    parser.add_argument("--collection", default=None)
    parser.add_argument("--hyde-count", type=int, default=0)
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()
    if not args.dir.exists():
        raise SystemExit(f"directory does not exist: {args.dir}")
    settings = get_settings()
    collection = args.collection or settings.qdrant_collection
    rows = build_rows(args.dir, max(0, min(3, args.hyde_count))) if args.dry_run else []
    manifest = args.dir / ".ingest_manifest.json"
    manifest.write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")
    if args.dry_run:
        print(json.dumps({"collection": collection, "chunks": len(rows), "manifest": str(manifest)}, ensure_ascii=False))
        return
    import asyncio
    from app.services.lightrag_service import get_lightrag_service

    text = "\n\n".join(page for path in sorted(args.dir.rglob("*")) if path.is_file() for _, page in _read_pages(path))
    doc_id = hashlib.sha256(str(args.dir.resolve()).encode("utf-8")).hexdigest()[:32]
    result = asyncio.run(get_lightrag_service().index_document(text, f"doc_{doc_id}", args.dir.name, args.collection or "default", 512))
    print(json.dumps({"collection": collection, "chunks": result.get("chunks", 0), "status": result.get("status"), "manifest": str(manifest)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
